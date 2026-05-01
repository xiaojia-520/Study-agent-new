from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "StudyAgent_中国计算机设计大赛_答辩PPT.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLOR = {
    "navy": RGBColor(22, 40, 74),
    "blue": RGBColor(54, 99, 216),
    "cyan": RGBColor(68, 176, 255),
    "teal": RGBColor(37, 156, 142),
    "green": RGBColor(27, 160, 98),
    "orange": RGBColor(241, 144, 42),
    "red": RGBColor(223, 78, 78),
    "bg": RGBColor(245, 247, 251),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(218, 224, 236),
    "text": RGBColor(32, 42, 58),
    "muted": RGBColor(104, 116, 136),
    "light": RGBColor(250, 251, 253),
}

FONT = "Microsoft YaHei"


def add_textbox(slide, x, y, w, h, text="", size=20, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLOR["text"]
    return box


def add_paragraph(frame, text, size=16, color=None, bold=False, level=0, bullet=False, align=PP_ALIGN.LEFT):
    p = frame.add_paragraph()
    p.alignment = align
    p.level = level
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLOR["text"]
    if bullet:
        p.text = f"• {text}"
    return p


def add_card(slide, x, y, w, h, fill=COLOR["card"], line=COLOR["line"], radius=0.08):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    shape.adjustments[0] = radius
    return shape


def add_badge(slide, x, y, w, h, text, fill, color=RGBColor(255, 255, 255), size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_header(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(7.8), Inches(0.5), title, size=28, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.82), Inches(9.4), Inches(0.3), subtitle, size=11, color=COLOR["muted"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.12), Inches(12.1), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR["blue"]
    line.line.color.rgb = COLOR["blue"]


def add_footer(slide, text="Study Agent | 中国计算机设计大赛答辩材料"):
    add_textbox(slide, Inches(0.58), Inches(7.08), Inches(4.6), Inches(0.2), text, size=9, color=COLOR["muted"])


def add_bullets_box(slide, x, y, w, h, title, bullets, accent):
    add_card(slide, x, y, w, h)
    add_badge(slide, x + Inches(0.18), y + Inches(0.16), Inches(1.05), Inches(0.28), title, accent)
    box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.56), w - Inches(0.36), h - Inches(0.72))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.name = FONT
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR["text"]
        else:
            add_paragraph(tf, item, size=15, bullet=True)
    return box


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["navy"])
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.45))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(17, 30, 56)
    band.line.color.rgb = RGBColor(17, 30, 56)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.85), Inches(13.333), Inches(0.65))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR["blue"]
    accent.line.color.rgb = COLOR["blue"]

    add_textbox(slide, Inches(0.75), Inches(1.15), Inches(5.6), Inches(0.4), "中国计算机设计大赛答辩", size=15, color=COLOR["cyan"], bold=True)
    add_textbox(slide, Inches(0.75), Inches(1.62), Inches(7.2), Inches(1.2), "Study Agent", size=30, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(0.78), Inches(2.45), Inches(8.4), Inches(0.8), "面向课堂学习场景的实时语音学习助手", size=22, color=RGBColor(236, 241, 248), bold=True)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.2),
        Inches(6.4),
        Inches(1.3),
        "核心能力：实时转写 / 课堂问答 / 多模态资料解析 / 历史回顾 / 视频字幕 / 课后学习闭环",
        size=15,
        color=RGBColor(215, 222, 236),
    )

    card = add_card(slide, Inches(8.2), Inches(1.25), Inches(4.25), Inches(4.85), fill=RGBColor(250, 252, 255), line=RGBColor(250, 252, 255))
    card.adjustments[0] = 0.04
    add_textbox(slide, Inches(8.55), Inches(1.55), Inches(2.3), Inches(0.35), "项目答辩信息", size=18, bold=True, color=COLOR["navy"])
    info = slide.shapes.add_textbox(Inches(8.55), Inches(2.02), Inches(3.4), Inches(3.55))
    tf = info.text_frame
    tf.clear()
    first = tf.paragraphs[0]
    first.text = "项目名称：Study Agent"
    first.font.name = FONT
    first.font.size = Pt(16)
    first.font.bold = True
    first.font.color.rgb = COLOR["text"]
    add_paragraph(tf, "项目定位：课堂场景下的学习智能体", size=15)
    add_paragraph(tf, "答辩方向：软件应用 / 智能教育", size=15)
    add_paragraph(tf, "学校 / 团队 / 指导教师：答辩前替换", size=15, color=COLOR["orange"], bold=True)
    add_paragraph(tf, "备注：PPT 已预留 StudyLibView 后续开发位置", size=14, color=COLOR["muted"])

    add_badge(slide, Inches(0.85), Inches(5.4), Inches(1.2), Inches(0.34), "本地优先", COLOR["green"])
    add_badge(slide, Inches(2.15), Inches(5.4), Inches(1.55), Inches(0.34), "实时转写 + RAG", COLOR["teal"])
    add_badge(slide, Inches(3.85), Inches(5.4), Inches(1.7), Inches(0.34), "多模态学习资料", COLOR["orange"])
    add_badge(slide, Inches(5.7), Inches(5.4), Inches(1.45), Inches(0.34), "后续可扩展", COLOR["blue"])
    return slide


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "1. 项目背景与痛点", "课堂学习的关键问题不在“获取内容”，而在“课中沉淀”和“课后复盘”")

    add_bullets_box(
        slide,
        Inches(0.65),
        Inches(1.55),
        Inches(3.9),
        Inches(4.7),
        "痛点一",
        [
            "课堂语速快、信息密度高，学生难以同时听课、记笔记、标重点。",
            "传统录音回放效率低，缺少结构化检索能力。",
        ],
        COLOR["blue"],
    )
    add_bullets_box(
        slide,
        Inches(4.72),
        Inches(1.55),
        Inches(3.9),
        Inches(4.7),
        "痛点二",
        [
            "课件、板书、图片、视频等资料分散，知识入口不统一。",
            "课后复习依赖人工整理，无法快速追问和定位知识点。",
        ],
        COLOR["orange"],
    )
    add_bullets_box(
        slide,
        Inches(8.79),
        Inches(1.55),
        Inches(3.9),
        Inches(4.7),
        "痛点三",
        [
            "已有学习工具常把课堂记录、知识检索、总结出题拆成多个系统。",
            "缺少从“上课”到“复习”再到“练习”的闭环。",
        ],
        COLOR["teal"],
    )
    add_footer(slide)
    return slide


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "2. 解决方案概述", "Study Agent 以课堂 session 为中心，把实时采集、知识组织和课后学习串成一条链路")

    add_badge(slide, Inches(0.75), Inches(1.45), Inches(1.0), Inches(0.32), "课中", COLOR["blue"])
    add_badge(slide, Inches(4.7), Inches(1.45), Inches(1.0), Inches(0.32), "课后", COLOR["orange"])
    add_badge(slide, Inches(8.72), Inches(1.45), Inches(1.45), Inches(0.32), "开发预留", COLOR["teal"])

    add_card(slide, Inches(0.7), Inches(1.88), Inches(3.65), Inches(4.6))
    add_textbox(slide, Inches(0.95), Inches(2.15), Inches(2.5), Inches(0.3), "课中智能辅助", size=20, bold=True, color=COLOR["navy"])
    box1 = slide.shapes.add_textbox(Inches(0.95), Inches(2.55), Inches(3.0), Inches(3.2))
    tf1 = box1.text_frame
    tf1.clear()
    first = tf1.paragraphs[0]
    first.text = "• 浏览器采集麦克风与课堂画面"
    first.font.name = FONT
    first.font.size = Pt(16)
    first.font.color.rgb = COLOR["text"]
    for item in [
        "• FunASR / Paraformer 实时转写",
        "• 课堂资料上传并自动解析",
        "• 视频录制、字幕生成、视觉帧识别",
        "• 实时指标反馈与会话管理",
    ]:
        add_paragraph(tf1, item, size=16)

    add_card(slide, Inches(4.82), Inches(1.88), Inches(3.65), Inches(4.6))
    add_textbox(slide, Inches(5.05), Inches(2.15), Inches(2.6), Inches(0.3), "课后知识利用", size=20, bold=True, color=COLOR["navy"])
    box2 = slide.shapes.add_textbox(Inches(5.05), Inches(2.55), Inches(3.0), Inches(3.2))
    tf2 = box2.text_frame
    tf2.clear()
    first2 = tf2.paragraphs[0]
    first2.text = "• 转写落地 SQLite + JSONL"
    first2.font.name = FONT
    first2.font.size = Pt(16)
    first2.font.color.rgb = COLOR["text"]
    for item in [
        "• Qdrant 准实时入库，支持课堂问答",
        "• 支持历史回顾、追问、转写精修",
        "• 基于课堂内容生成总结与测验",
    ]:
        add_paragraph(tf2, item, size=16)

    add_card(slide, Inches(8.94), Inches(1.88), Inches(3.65), Inches(4.6), fill=RGBColor(247, 252, 251), line=RGBColor(177, 223, 214))
    add_textbox(slide, Inches(9.18), Inches(2.15), Inches(3.0), Inches(0.3), "StudyLibView 后续扩展", size=20, bold=True, color=COLOR["navy"])
    box3 = slide.shapes.add_textbox(Inches(9.18), Inches(2.55), Inches(3.0), Inches(3.35))
    tf3 = box3.text_frame
    tf3.clear()
    first3 = tf3.paragraphs[0]
    first3.text = "• 已预留 /workshop 页面入口"
    first3.font.name = FONT
    first3.font.size = Pt(16)
    first3.font.color.rgb = COLOR["text"]
    for item in [
        "• 自动总结出笔记",
        "• 课后出题与自测",
        "• 为答辩演示保留独立功能页",
        "• 当前不在已完成功能内，避免虚报进度",
    ]:
        add_paragraph(tf3, item, size=16)
    add_footer(slide)
    return slide


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "3. 现阶段功能全景", "系统围绕一堂课的完整学习流程组织功能")

    titles = [
        ("实时语音转写", "麦克风音频通过 WebSocket 进入后端，输出 partial / final transcript。", COLOR["blue"]),
        ("课堂问答", "按当前课次、课程历史或全局范围检索，支持 LLM 直接回答。", COLOR["teal"]),
        ("课件资料入库", "PDF / PPT / 图片 / HTML 通过 MinerU 解析后进入知识库。", COLOR["orange"]),
        ("历史回顾", "查看课程、转写、问答、精修结果和视频字幕。", COLOR["green"]),
        ("视频字幕", "课堂录像可转为字幕与文本，再反向补充知识库。", COLOR["blue"]),
        ("学习生成", "已有后端能力支持总结与测验，前端学习工坊页待接入。", COLOR["teal"]),
    ]

    positions = [
        (Inches(0.72), Inches(1.6)),
        (Inches(4.45), Inches(1.6)),
        (Inches(8.18), Inches(1.6)),
        (Inches(0.72), Inches(4.0)),
        (Inches(4.45), Inches(4.0)),
        (Inches(8.18), Inches(4.0)),
    ]

    for (title, desc, accent), (x, y) in zip(titles, positions):
        add_card(slide, x, y, Inches(3.45), Inches(2.05))
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(3.45), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.color.rgb = accent
        add_textbox(slide, x + Inches(0.16), y + Inches(0.2), Inches(2.8), Inches(0.3), title, size=18, bold=True, color=COLOR["navy"])
        add_textbox(slide, x + Inches(0.16), y + Inches(0.62), Inches(3.0), Inches(1.05), desc, size=14, color=COLOR["text"])
    add_footer(slide)
    return slide


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "4. 系统架构", "前后端解耦，数据存储与模型运行分层清晰，便于本地部署与功能扩展")

    cols = [
        (Inches(0.8), "前端层", COLOR["blue"], ["Vue 3", "TypeScript", "Pinia", "Vite", "Live / History / Workshop"]),
        (Inches(4.45), "服务层", COLOR["teal"], ["FastAPI", "WebSocket", "Session API", "History API", "Asset / Video / Vision API"]),
        (Inches(8.1), "智能层与存储层", COLOR["orange"], ["FunASR / Paraformer", "BGE Embedding", "Qdrant", "SQLite + JSONL", "MinerU / OpenAI Compatible LLM"]),
    ]

    for x, title, accent, items in cols:
        add_card(slide, x, Inches(1.7), Inches(3.0), Inches(4.8))
        add_badge(slide, x + Inches(0.18), Inches(1.88), Inches(1.5), Inches(0.3), title, accent)
        box = slide.shapes.add_textbox(x + Inches(0.18), Inches(2.35), Inches(2.55), Inches(3.8))
        tf = box.text_frame
        tf.clear()
        first = tf.paragraphs[0]
        first.text = f"• {items[0]}"
        first.font.name = FONT
        first.font.size = Pt(16)
        first.font.color.rgb = COLOR["text"]
        for item in items[1:]:
            add_paragraph(tf, item, size=16, bullet=True)

    for x1, x2 in [(Inches(3.85), Inches(4.35)), (Inches(7.5), Inches(8.0))]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, Inches(4.1), x2, Inches(4.1))
        conn.line.color.rgb = COLOR["blue"]
        conn.line.width = Pt(2)
    add_footer(slide)
    return slide


def slide_realtime_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "5. 实时语音链路", "课堂语音从浏览器进入后端流水线，再回到界面并落地存储")

    steps = [
        ("01 浏览器采集", "麦克风音频 -> AudioContext -> float32 音频帧"),
        ("02 WebSocket 传输", "按 session 建立连接，连续推送音频与心跳"),
        ("03 VAD + ASR", "FrameSlicer + VADProcessor + Realtime ASR Driver"),
        ("04 实时回显", "前端显示 partial / final transcript 与音频指标"),
        ("05 数据落地", "SQLite / JSONL 持久化，并准备进入后续 RAG"),
    ]

    x = Inches(0.7)
    for idx, (title, desc) in enumerate(steps):
        width = Inches(2.35)
        add_card(slide, x, Inches(2.3), width, Inches(2.1))
        add_badge(slide, x + Inches(0.14), Inches(2.48), Inches(0.82), Inches(0.28), f"{idx + 1}", COLOR["blue"])
        add_textbox(slide, x + Inches(0.14), Inches(2.85), width - Inches(0.28), Inches(0.34), title, size=17, bold=True, color=COLOR["navy"])
        add_textbox(slide, x + Inches(0.14), Inches(3.28), width - Inches(0.28), Inches(0.82), desc, size=13)
        if idx < len(steps) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + width, Inches(3.35), x + width + Inches(0.35), Inches(3.35))
            conn.line.color.rgb = COLOR["teal"]
            conn.line.width = Pt(2)
        x += Inches(2.52)
    add_footer(slide)
    return slide


def slide_rag_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "6. RAG 与知识组织", "本项目不是简单聊天，而是围绕课堂 session 构建可检索、可追问、可沉淀的知识流")

    left = add_card(slide, Inches(0.78), Inches(1.65), Inches(3.55), Inches(4.95))
    add_textbox(slide, Inches(1.02), Inches(1.95), Inches(2.6), Inches(0.35), "知识来源", size=19, bold=True, color=COLOR["navy"])
    left_box = slide.shapes.add_textbox(Inches(1.02), Inches(2.35), Inches(2.95), Inches(3.7))
    tf = left_box.text_frame
    tf.clear()
    first = tf.paragraphs[0]
    first.text = "• 实时语音最终转写"
    first.font.name = FONT
    first.font.size = Pt(16)
    first.font.color.rgb = COLOR["text"]
    for item in [
        "• 课件 / 文档 / HTML / 图片解析结果",
        "• 视频字幕与视觉区域识别文本",
        "• 历史 session 的可追溯课堂内容",
    ]:
        add_paragraph(tf, item, size=16, bullet=True)

    center = add_card(slide, Inches(4.88), Inches(1.65), Inches(3.55), Inches(4.95))
    add_textbox(slide, Inches(5.1), Inches(1.95), Inches(2.6), Inches(0.35), "组织方式", size=19, bold=True, color=COLOR["navy"])
    center_box = slide.shapes.add_textbox(Inches(5.1), Inches(2.35), Inches(2.95), Inches(3.7))
    tf2 = center_box.text_frame
    tf2.clear()
    first2 = tf2.paragraphs[0]
    first2.text = "• 统一转为 TranscriptRecord / Chunk"
    first2.font.name = FONT
    first2.font.size = Pt(16)
    first2.font.color.rgb = COLOR["text"]
    for item in [
        "• 附带 course_id / lesson_id / source_type 等元数据",
        "• 按阈值缓冲写入 Qdrant",
        "• 支持当前课次、课程历史、全局范围检索",
    ]:
        add_paragraph(tf2, item, size=16, bullet=True)

    right = add_card(slide, Inches(8.98), Inches(1.65), Inches(3.55), Inches(4.95))
    add_textbox(slide, Inches(9.2), Inches(1.95), Inches(2.6), Inches(0.35), "结果利用", size=19, bold=True, color=COLOR["navy"])
    right_box = slide.shapes.add_textbox(Inches(9.2), Inches(2.35), Inches(2.95), Inches(3.7))
    tf3 = right_box.text_frame
    tf3.clear()
    first3 = tf3.paragraphs[0]
    first3.text = "• 实时问答与历史追问"
    first3.font.name = FONT
    first3.font.size = Pt(16)
    first3.font.color.rgb = COLOR["text"]
    for item in [
        "• 课堂总结、知识点提炼",
        "• 课后测验生成",
        "• 后续 StudyLibView 笔记页与练习页",
    ]:
        add_paragraph(tf3, item, size=16, bullet=True)

    for x1, x2 in [(Inches(4.33), Inches(4.83)), (Inches(8.43), Inches(8.93))]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, Inches(4.1), x2, Inches(4.1))
        conn.line.color.rgb = COLOR["orange"]
        conn.line.width = Pt(2)
    add_footer(slide)
    return slide


def slide_ui(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "7. 界面与交互组织", "当前前端已形成较清晰的三栏工作台，并预留后续学习工坊入口")

    add_card(slide, Inches(0.72), Inches(1.7), Inches(7.95), Inches(4.95), fill=RGBColor(248, 250, 255))
    add_textbox(slide, Inches(0.94), Inches(1.95), Inches(2.6), Inches(0.3), "主界面结构示意", size=19, bold=True, color=COLOR["navy"])

    left = add_card(slide, Inches(1.0), Inches(2.55), Inches(2.08), Inches(3.35), fill=RGBColor(255, 255, 255))
    mid = add_card(slide, Inches(3.28), Inches(2.55), Inches(2.2), Inches(3.35), fill=RGBColor(255, 255, 255))
    right_top = add_card(slide, Inches(5.68), Inches(2.55), Inches(2.55), Inches(1.45), fill=RGBColor(255, 255, 255))
    right_bottom = add_card(slide, Inches(5.68), Inches(4.18), Inches(2.55), Inches(1.72), fill=RGBColor(255, 255, 255))

    add_textbox(slide, Inches(1.16), Inches(2.85), Inches(1.7), Inches(0.25), "Session / 设备 / 资料", size=16, bold=True)
    add_textbox(slide, Inches(1.16), Inches(3.2), Inches(1.6), Inches(2.1), "课程名\n模型选择\n麦克风 / 摄像头\n素材上传状态", size=14, color=COLOR["muted"])

    add_textbox(slide, Inches(3.44), Inches(2.85), Inches(1.7), Inches(0.25), "实时转写面板", size=16, bold=True)
    add_textbox(slide, Inches(3.44), Inches(3.2), Inches(1.7), Inches(2.1), "按时间顺序显示\npartial / final transcript\n自动滚动与同步", size=14, color=COLOR["muted"])

    add_textbox(slide, Inches(5.86), Inches(2.82), Inches(1.7), Inches(0.25), "视频 / 视觉区", size=16, bold=True)
    add_textbox(slide, Inches(5.86), Inches(3.16), Inches(1.95), Inches(0.48), "摄像头预览、区域框选、视频录制与字幕状态", size=13, color=COLOR["muted"])

    add_textbox(slide, Inches(5.86), Inches(4.46), Inches(1.7), Inches(0.25), "课堂问答", size=16, bold=True)
    add_textbox(slide, Inches(5.86), Inches(4.82), Inches(1.95), Inches(0.62), "当前 session 问答、RAG 开关、追问与结果回显", size=13, color=COLOR["muted"])

    add_card(slide, Inches(9.0), Inches(1.7), Inches(3.55), Inches(4.95))
    add_textbox(slide, Inches(9.22), Inches(1.95), Inches(2.6), Inches(0.3), "页面路由现状", size=19, bold=True, color=COLOR["navy"])
    box = slide.shapes.add_textbox(Inches(9.22), Inches(2.35), Inches(2.95), Inches(3.75))
    tf = box.text_frame
    tf.clear()
    first = tf.paragraphs[0]
    first.text = "• /          实时课堂主界面"
    first.font.name = FONT
    first.font.size = Pt(16)
    first.font.color.rgb = COLOR["text"]
    for item in [
        "• /history   历史回顾页",
        "• /workshop  StudyLibView 预留入口",
        "• 当前 workshop 页仅占位，适合继续扩展答辩功能",
    ]:
        add_paragraph(tf, item, size=16, bullet=True)
    add_footer(slide)
    return slide


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "8. 答辩演示路径", "建议现场按“课中 -> 课后 -> 规划扩展”的顺序演示，逻辑最顺")

    steps = [
        ("创建课堂会话", "输入课程名，选择模型与设备，创建 session。"),
        ("开始实时录音", "展示 partial / final transcript 的实时回显。"),
        ("上传课堂资料", "导入 PPT / PDF / 图片，说明解析入库逻辑。"),
        ("课堂追问", "围绕当前课程提问，展示历史上下文与检索结果。"),
        ("历史回顾", "切换 History 页面，展示转写、追问、字幕与精修记录。"),
        ("StudyLibView 预告", "说明后续会把笔记总结和课后出题统一放在 /workshop。"),
    ]

    y = Inches(1.58)
    for idx, (title, desc) in enumerate(steps, start=1):
        add_badge(slide, Inches(0.9), y + Inches(0.06), Inches(0.58), Inches(0.3), str(idx), COLOR["blue"])
        add_card(slide, Inches(1.65), y, Inches(10.75), Inches(0.72))
        add_textbox(slide, Inches(1.95), y + Inches(0.12), Inches(2.3), Inches(0.22), title, size=17, bold=True, color=COLOR["navy"])
        add_textbox(slide, Inches(4.15), y + Inches(0.12), Inches(7.6), Inches(0.3), desc, size=14, color=COLOR["text"])
        y += Inches(0.82)
    add_footer(slide)
    return slide


def slide_status(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "9. 完成度与比赛表述边界", "答辩时应区分“已完成”“已具备后端能力”“前端预留待接入”，避免把规划说成现状")

    rows = [
        ("实时语音转写", "已完成", COLOR["green"], "主链路已打通，可展示实时输出与落地。"),
        ("课堂资料解析入库", "已完成", COLOR["green"], "课件 / 图片 / HTML / 视频字幕已进入统一知识流。"),
        ("课堂问答与历史回顾", "已完成", COLOR["green"], "支持当前 session 问答、历史页面查看与追问。"),
        ("课堂总结 / 课后出题后端能力", "已具备", COLOR["orange"], "后端已有 summary / quiz service，可作为能力说明。"),
        ("StudyLibView 前端页面", "已预留", COLOR["orange"], "当前只有 /workshop 路由占位，适合继续开发答辩展示页。"),
        ("自动总结出笔记 / 课后出题前端闭环", "待实现", COLOR["red"], "建议答辩中明确说明为下一步开发重点。"),
    ]

    y = Inches(1.55)
    for feature, status, status_color, note in rows:
        add_card(slide, Inches(0.7), y, Inches(12.0), Inches(0.72))
        add_textbox(slide, Inches(0.95), y + Inches(0.15), Inches(3.3), Inches(0.24), feature, size=15, bold=True)
        add_badge(slide, Inches(4.55), y + Inches(0.17), Inches(1.02), Inches(0.28), status, status_color)
        add_textbox(slide, Inches(5.85), y + Inches(0.15), Inches(6.45), Inches(0.24), note, size=14, color=COLOR["muted"])
        y += Inches(0.82)
    add_footer(slide)
    return slide


def slide_workshop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "10. StudyLibView 开发预留页", "本页专门为你后续实现两个功能预留，答辩时可以作为“下一阶段重点”或“扩展规划”展示")

    add_badge(slide, Inches(0.84), Inches(1.45), Inches(2.0), Inches(0.32), "入口：/workshop", COLOR["teal"])

    left = add_card(slide, Inches(0.72), Inches(1.95), Inches(5.9), Inches(4.95), fill=RGBColor(249, 252, 255))
    right = add_card(slide, Inches(6.76), Inches(1.95), Inches(5.85), Inches(4.95), fill=RGBColor(249, 252, 255))

    add_textbox(slide, Inches(0.98), Inches(2.18), Inches(3.5), Inches(0.3), "功能一：自动总结出笔记", size=21, bold=True, color=COLOR["navy"])
    left_text = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(4.0), Inches(2.6))
    tf1 = left_text.text_frame
    tf1.clear()
    first = tf1.paragraphs[0]
    first.text = "• 输入：课堂转写 + 课件解析 + 视频字幕"
    first.font.name = FONT
    first.font.size = Pt(16)
    first.font.color.rgb = COLOR["text"]
    for item in [
        "• 输出：结构化笔记、重点概念、复习清单",
        "• 形式：一键生成，可沉淀为课后学习页",
        "• 建议展示：笔记卡片、时间线、章节摘要",
    ]:
        add_paragraph(tf1, item, size=16, bullet=True)
    placeholder1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(4.95), Inches(1.28))
    placeholder1.fill.solid()
    placeholder1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    placeholder1.line.color.rgb = COLOR["line"]
    tfp1 = placeholder1.text_frame
    tfp1.text = "此处替换为 StudyLibView 笔记页截图 / 演示稿"
    tfp1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tfp1.paragraphs[0].font.name = FONT
    tfp1.paragraphs[0].font.size = Pt(14)
    tfp1.paragraphs[0].font.color.rgb = COLOR["muted"]
    tfp1.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, Inches(7.02), Inches(2.18), Inches(3.5), Inches(0.3), "功能二：课后出题与自测", size=21, bold=True, color=COLOR["navy"])
    right_text = slide.shapes.add_textbox(Inches(7.04), Inches(2.6), Inches(4.0), Inches(2.6))
    tf2 = right_text.text_frame
    tf2.clear()
    first2 = tf2.paragraphs[0]
    first2.text = "• 输入：课堂知识点与总结结果"
    first2.font.name = FONT
    first2.font.size = Pt(16)
    first2.font.color.rgb = COLOR["text"]
    for item in [
        "• 输出：选择题 / 简答题 / 重点回忆题",
        "• 形式：即时练习 + 答案解析 + 知识回溯",
        "• 建议展示：题目卡片、作答反馈、错题回看",
    ]:
        add_paragraph(tf2, item, size=16, bullet=True)
    placeholder2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.04), Inches(5.2), Inches(4.9), Inches(1.28))
    placeholder2.fill.solid()
    placeholder2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    placeholder2.line.color.rgb = COLOR["line"]
    tfp2 = placeholder2.text_frame
    tfp2.text = "此处替换为 StudyLibView 出题页截图 / 演示稿"
    tfp2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tfp2.paragraphs[0].font.name = FONT
    tfp2.paragraphs[0].font.size = Pt(14)
    tfp2.paragraphs[0].font.color.rgb = COLOR["muted"]
    tfp2.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_footer(slide)
    return slide


def slide_innovation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "11. 项目亮点与创新点", "本项目的价值不只是识别准确，而是把课堂资料变成可以持续利用的学习资产")

    items = [
        ("以课堂 session 为主轴", "不是单纯录音或聊天，而是把课程、课次、资料、问答和历史记录绑定在一起。", COLOR["blue"]),
        ("多源知识统一进入 RAG", "实时转写、课件解析、视频字幕、视觉文本统一带元数据进入知识库。", COLOR["orange"]),
        ("本地优先，适合校园场景", "Embedding、Qdrant、ASR 模型支持本地运行，便于离线演示和数据留存。", COLOR["green"]),
        ("具备课后学习延展性", "后端已经具备总结和出题服务，前端可继续扩为 StudyLibView 学习工坊。", COLOR["teal"]),
    ]

    y = Inches(1.68)
    for title, desc, accent in items:
        add_card(slide, Inches(0.78), y, Inches(11.8), Inches(0.9))
        accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.78), y, Inches(0.12), Inches(0.9))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.color.rgb = accent
        add_textbox(slide, Inches(1.08), y + Inches(0.14), Inches(3.2), Inches(0.24), title, size=16, bold=True, color=COLOR["navy"])
        add_textbox(slide, Inches(4.22), y + Inches(0.14), Inches(7.85), Inches(0.34), desc, size=14, color=COLOR["text"])
        y += Inches(1.05)
    add_footer(slide)
    return slide


def slide_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["bg"])
    add_header(slide, "12. 后续开发计划", "比赛前后可以围绕 StudyLibView 把“能力存在”升级为“完整学习体验”")

    phases = [
        ("短期", "接入 StudyLibView：自动笔记、课后出题、页面交互闭环", COLOR["blue"]),
        ("中期", "增加错题本、知识卡片、复习日历、章节知识树", COLOR["orange"]),
        ("长期", "引入学习画像、个性化推荐、跨课程知识关联", COLOR["teal"]),
    ]

    x = Inches(0.95)
    for title, desc, accent in phases:
        add_card(slide, x, Inches(2.18), Inches(3.75), Inches(2.85))
        add_badge(slide, x + Inches(0.18), Inches(2.38), Inches(0.95), Inches(0.28), title, accent)
        add_textbox(slide, x + Inches(0.18), Inches(2.86), Inches(3.25), Inches(1.2), desc, size=17, bold=True, color=COLOR["navy"])
        x += Inches(4.03)

    add_card(slide, Inches(1.05), Inches(5.45), Inches(11.2), Inches(0.82), fill=RGBColor(247, 251, 255))
    add_textbox(
        slide,
        Inches(1.25),
        Inches(5.68),
        Inches(10.7),
        Inches(0.26),
        "建议答辩措辞：当前系统已完成课堂主链路，StudyLibView 将把现有总结/出题能力进一步产品化、页面化、可演示化。",
        size=15,
        color=COLOR["navy"],
        bold=True,
    )
    add_footer(slide)
    return slide


def slide_ending(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR["navy"])
    add_textbox(slide, Inches(0.95), Inches(1.55), Inches(5.4), Inches(0.4), "答辩总结", size=18, color=COLOR["cyan"], bold=True)
    add_textbox(slide, Inches(0.95), Inches(2.1), Inches(6.8), Inches(1.0), "Study Agent", size=30, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(
        slide,
        Inches(0.98),
        Inches(3.0),
        Inches(7.2),
        Inches(1.5),
        "我们尝试把课堂上的声音、画面、资料和课后学习动作连接起来，\n把“听完一节课”升级为“沉淀一份可持续利用的学习资产”。",
        size=20,
        color=RGBColor(232, 239, 250),
    )
    add_badge(slide, Inches(0.98), Inches(5.15), Inches(1.4), Inches(0.34), "实时课堂", COLOR["blue"])
    add_badge(slide, Inches(2.55), Inches(5.15), Inches(1.65), Inches(0.34), "知识沉淀", COLOR["teal"])
    add_badge(slide, Inches(4.38), Inches(5.15), Inches(1.55), Inches(0.34), "课后学习", COLOR["orange"])
    add_badge(slide, Inches(6.1), Inches(5.15), Inches(1.95), Inches(0.34), "持续扩展能力", COLOR["green"])
    add_textbox(slide, Inches(9.15), Inches(2.2), Inches(3.0), Inches(0.7), "感谢评委老师聆听", size=24, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.25), Inches(3.15), Inches(2.8), Inches(0.7), "Q & A", size=34, bold=True, color=COLOR["cyan"], align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.95), Inches(4.4), Inches(3.3), Inches(1.2), "提示：答辩前请补齐学校、团队成员、指导老师等信息。", size=14, color=RGBColor(214, 224, 238), align=PP_ALIGN.CENTER)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_features(prs)
    slide_architecture(prs)
    slide_realtime_flow(prs)
    slide_rag_flow(prs)
    slide_ui(prs)
    slide_demo(prs)
    slide_status(prs)
    slide_workshop(prs)
    slide_innovation(prs)
    slide_plan(prs)
    slide_ending(prs)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
